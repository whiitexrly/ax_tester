"""Script for automating creation of .pptx file for presentation and demo session.

This file is fully built in vibe-coding.
"""

import json
import mimetypes
import re
from collections import Counter
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib.parse import urlparse
from urllib.request import Request, urlopen

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from utils.wcag_helper import WCAG_RULE_MAPPER, get_wcag_level

ASSETS_DIR = Path(__file__).resolve().parents[1] / "assets"
TEMPLATES_DIR = ASSETS_DIR / "templates"
LEVEL_ORDER = {"A": 0, "AA": 1, "AAA": 2}
MAX_WCAG_ROWS_PER_SLIDE = 13
REPLY_GREEN = (0, 177, 63)
STATUS_OK_GREEN = (0, 177, 63)
STATUS_ERROR_RED = (211, 47, 47)
STATUS_OK_SYMBOL = "✓"
STATUS_ERROR_SYMBOL = "✗"
FIRSTSLIDE_TEMPLATE_FILENAME = "template_firstslide.pdf"
THANKYOU_TEMPLATE_FILENAME = "template_thankyou.pdf"
TEMPLATE_LOGO_FILENAME = "template_logo.png"
REPORT_FILENAME = "ax_report.json"
LOGO_WIDTH_IN = 0.68
LOGO_MARGIN_IN = 0.16
SLIDE_WIDTH_16_9_IN = 13.333
SLIDE_HEIGHT_16_9_IN = 7.5
FONT_SIZE_SCALE = 1.12
BODY_FONT_NAME = "Arial"
TITLE_FONT_NAME = "Arial Black"
SNIPPET_FONT_NAME = "Consolas"
OUTCOME_FULLY_COMPLIANT_MAX_ERRORS = 5
OUTCOME_PARTIALLY_COMPLIANT_MAX_ERRORS = 50
OUTCOME_FULLY_COMPLIANT = "Fully Compliant"
OUTCOME_PARTIALLY_COMPLIANT = "Partially Compliant"
OUTCOME_NOT_COMPLIANT = "Not Compliant"
OUTCOME_ORANGE = (245, 124, 0)
ISSUE_IMAGE_DOWNLOAD_TIMEOUT = 4
ISSUE_IMAGE_MAX_BYTES = 5_000_000
HTML_SNIPPET_MAX_CHARS = 500
WHY_IT_MATTERS_MAX_CHARS = 260
POTENTIAL_EXPOSURES_MAX_CHARS = 420
FIX_MAX_CHARS = 340


def build_pptx_report(results_dir: str) -> str:
    if not results_dir:
        raise ValueError("results_dir not found in tool_context.state. Run save first.")

    report = _load_report(results_dir)
    all_issues = _load_all_issues(report)
    all_issues = _sort_issues(all_issues)
    wcag_summary = _build_wcag_summary(all_issues)
    overview = _build_overview(report, all_issues, wcag_summary)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_16_9_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_16_9_IN)
    _add_firstslide_template_slide(prs, RGBColor, Inches, Pt)
    _add_overview_slide(prs, overview, RGBColor, Inches, Pt)
    _add_wcag_summary_slides(prs, wcag_summary, RGBColor, Inches, Pt)

    issue_image_cache: dict[str, bytes | None] = {}
    for issue_index, issue in enumerate(all_issues, start=1):
        _add_issue_slide(
            prs,
            issue_index,
            len(all_issues),
            issue,
            RGBColor,
            MSO_ANCHOR,
            MSO_AUTO_SIZE,
            Inches,
            Pt,
            issue_image_cache,
        )

    _add_thankyou_template_slide(prs, RGBColor, Inches, Pt)
    _add_logo_to_middle_slides(prs, Inches)

    out_pptx = str(Path(results_dir) / "ax_report.pptx")
    prs.save(out_pptx)
    return out_pptx


def _load_report(results_dir: str) -> dict[str, Any]:
    report_path = Path(results_dir) / REPORT_FILENAME
    if not report_path.exists():
        raise FileNotFoundError(f"Missing report file: {report_path}")

    with open(report_path, encoding="utf-8") as file:
        data = json.load(file)
    if not isinstance(data, dict):
        raise ValueError(f"Invalid report format in {report_path}: expected JSON object")
    return data


def _load_all_issues(report: dict[str, Any]) -> list[dict[str, Any]]:
    issue_list = report.get("issue_list", [])
    if not isinstance(issue_list, list):
        return []
    return [issue.copy() for issue in issue_list if isinstance(issue, dict)]


def _sort_issues(all_issues: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return sorted(
        all_issues,
        key=lambda issue: (
            _level_order(get_wcag_level(issue.get("wcag_rule"))),
            str(issue.get("wcag_rule", "")),
            str(issue.get("source", "")),
            str(issue.get("id", "")),
        ),
    )


def _build_wcag_summary(all_issues: list[dict[str, Any]]) -> list[dict[str, Any]]:
    rule_counts = Counter(
        issue.get("wcag_rule")
        for issue in all_issues
        if isinstance(issue.get("wcag_rule"), str) and issue.get("wcag_rule")
    )

    baseline_rules = _get_a_aa_wcag_rules()
    extra_rules = sorted(rule for rule in rule_counts if rule not in baseline_rules)
    all_rules = [*baseline_rules, *extra_rules]

    rows: list[dict[str, Any]] = []
    for wcag_rule in all_rules:
        total_issues = int(rule_counts.get(wcag_rule, 0))
        rows.append(
            {
                "wcag_rule": wcag_rule,
                "level": get_wcag_level(wcag_rule),
                "total_issues": total_issues,
                "status": STATUS_OK_SYMBOL if total_issues == 0 else STATUS_ERROR_SYMBOL,
            }
        )

    return sorted(
        rows,
        key=lambda row: (
            _level_order(row.get("level")),
            -int(row.get("total_issues", 0)),
            str(row.get("wcag_rule", "")),
        ),
    )


def _build_overview(
    report: dict[str, Any],
    all_issues: list[dict[str, Any]],
    wcag_summary: list[dict[str, Any]],
) -> dict[str, Any]:
    level_counts = {"A": 0, "AA": 0, "AAA": 0}
    for issue in all_issues:
        level_counts[get_wcag_level(issue.get("wcag_rule"))] += 1

    zero_issue_rules = sum(1 for row in wcag_summary if int(row.get("total_issues", 0)) == 0)
    total_rules = len(wcag_summary)
    violated_rules = total_rules - zero_issue_rules

    report_tool_name = str(report.get("tool_name", "")).strip()
    page_url = str(report.get("page", "")).strip()
    final_outcome, final_outcome_color = _get_final_outcome(int(len(all_issues)))

    return {
        "tool_name": report_tool_name if report_tool_name else "AX Tester",
        "url": page_url if page_url else "N/A",
        "total_issues": len(all_issues),
        "level_counts": level_counts,
        "violated_rules": violated_rules,
        "total_rules": total_rules,
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "final_outcome": final_outcome,
        "final_outcome_color": final_outcome_color,
    }


def _add_overview_slide(prs: Any, overview: dict[str, Any], rgb_color: Any, inches: Any, pt: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(inches(0.5), inches(0.2), inches(12.2), inches(0.6))
    title_frame = title_shape.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = "Accessibility Report Summary"
    title_paragraph.font.size = _scaled_pt(pt, 30)
    title_paragraph.font.bold = True
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.color.rgb = _to_rgb(rgb_color, REPLY_GREEN)

    details_shape = slide.shapes.add_textbox(inches(0.6), inches(1.0), inches(6.95), inches(5.8))
    details_frame = details_shape.text_frame
    details_frame.clear()

    lines = [
        ("Tool", overview["tool_name"]),
        ("URL analyzed", overview["url"]),
        ("Total issues", overview["total_issues"]),
        ("Rules violated", f"{overview['violated_rules']} / {overview['total_rules']}"),
        ("Final outcome", overview["final_outcome"]),
    ]

    for idx, (label, value) in enumerate(lines):
        paragraph = details_frame.paragraphs[0] if idx == 0 else details_frame.add_paragraph()
        if label == "Final outcome":
            _set_label_value_paragraph(
                paragraph=paragraph,
                label=str(label),
                value=str(value),
                pt=pt,
                font_size=17,
                value_bold=True,
                value_color=overview["final_outcome_color"],
                rgb_color=rgb_color,
            )
        else:
            _set_label_value_paragraph(paragraph, str(label), str(value), pt, 17)
        paragraph.level = 0

    wcag_table_shape = slide.shapes.add_table(4, 2, inches(7.9), inches(1.05), inches(4.85), inches(2.55))
    wcag_table = wcag_table_shape.table
    wcag_table.columns[0].width = inches(2.6)
    wcag_table.columns[1].width = inches(2.25)

    wcag_table.cell(0, 0).text = "WCAG Level"
    wcag_table.cell(0, 1).text = "Bug Count"
    for col_index in [0, 1]:
        cell = wcag_table.cell(0, col_index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _to_rgb(rgb_color, REPLY_GREEN)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.name = BODY_FONT_NAME
        paragraph.font.size = _scaled_pt(pt, 13)
        paragraph.font.color.rgb = rgb_color(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    wcag_rows = [
        ("A", overview["level_counts"]["A"]),
        ("AA", overview["level_counts"]["AA"]),
        ("AAA", overview["level_counts"]["AAA"]),
    ]
    for row_index, (level_label, bug_count) in enumerate(wcag_rows, start=1):
        wcag_table.cell(row_index, 0).text = str(level_label)
        wcag_table.cell(row_index, 1).text = str(bug_count)
        level_paragraph = wcag_table.cell(row_index, 0).text_frame.paragraphs[0]
        level_paragraph.font.bold = True
        level_paragraph.font.name = BODY_FONT_NAME
        level_paragraph.font.size = _scaled_pt(pt, 12)
        level_paragraph.alignment = PP_ALIGN.CENTER
        count_paragraph = wcag_table.cell(row_index, 1).text_frame.paragraphs[0]
        count_paragraph.font.bold = True
        count_paragraph.font.name = BODY_FONT_NAME
        count_paragraph.font.size = _scaled_pt(pt, 12)
        count_paragraph.alignment = PP_ALIGN.CENTER

    generated_shape = slide.shapes.add_textbox(inches(0.6), inches(6.95), inches(6.95), inches(0.3))
    generated_frame = generated_shape.text_frame
    generated_frame.clear()
    generated_paragraph = generated_frame.paragraphs[0]
    _set_label_value_paragraph(generated_paragraph, "Generated at", str(overview["generated_at"]), pt, 15)


def _add_wcag_summary_slides(
    prs: Any,
    wcag_summary: list[dict[str, Any]],
    rgb_color: Any,
    inches: Any,
    pt: Any,
) -> None:
    chunks = _chunk_rows(wcag_summary, MAX_WCAG_ROWS_PER_SLIDE)
    total_chunks = len(chunks)
    for chunk_index, chunk in enumerate(chunks, start=1):
        _add_wcag_summary_slide(prs, chunk, chunk_index, total_chunks, rgb_color, inches, pt)


def _add_wcag_summary_slide(
    prs: Any,
    wcag_summary: list[dict[str, Any]],
    chunk_index: int,
    total_chunks: int,
    rgb_color: Any,
    inches: Any,
    pt: Any,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(inches(0.3), inches(0.15), inches(9.4), inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    if total_chunks > 1:
        title_paragraph.text = f"WCAG Rules Summary {chunk_index}/{total_chunks}"
    else:
        title_paragraph.text = "WCAG Rules Summary"
    title_paragraph.font.size = _scaled_pt(pt, 22)
    title_paragraph.font.bold = True
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.color.rgb = _to_rgb(rgb_color, REPLY_GREEN)

    row_count = max(len(wcag_summary), 1) + 1
    table_shape = slide.shapes.add_table(row_count, 4, inches(0.25), inches(0.75), inches(11.8), inches(6.5))
    table = table_shape.table

    table.columns[0].width = inches(7.5)
    table.columns[1].width = inches(1.1)
    table.columns[2].width = inches(1.3)
    table.columns[3].width = inches(1.5)

    headers = ["WCAG Rule", "Level", "Errors", "Status"]
    for col_index, header_text in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = _to_rgb(rgb_color, REPLY_GREEN)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = _scaled_pt(pt, 13)
        paragraph.font.bold = True
        paragraph.font.name = BODY_FONT_NAME
        paragraph.font.color.rgb = rgb_color(255, 255, 255)

    if not wcag_summary:
        table.cell(1, 0).text = "No WCAG rules available."
        for col_index in [1, 2, 3]:
            table.cell(1, col_index).text = "-"
        for col_index in [0, 1, 2, 3]:
            paragraph = table.cell(1, col_index).text_frame.paragraphs[0]
            paragraph.font.name = BODY_FONT_NAME
            paragraph.font.size = _scaled_pt(pt, 11)
        return

    body_font_size = 11
    for row_index, row in enumerate(wcag_summary, start=1):
        table.cell(row_index, 0).text = _display_wcag_rule(str(row.get("wcag_rule", "N/A")))
        table.cell(row_index, 1).text = str(row.get("level", "AAA"))
        table.cell(row_index, 2).text = str(row.get("total_issues", 0))
        table.cell(row_index, 3).text = str(row.get("status", STATUS_ERROR_SYMBOL))

        for col_index in [0, 1, 2, 3]:
            cell = table.cell(row_index, col_index)
            paragraph = cell.text_frame.paragraphs[0]
            font_size = 14 if col_index == 0 else body_font_size
            paragraph.font.size = _scaled_pt(pt, font_size)
            paragraph.font.name = BODY_FONT_NAME
            paragraph.font.bold = col_index == 3

            if col_index == 3:
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                paragraph.alignment = PP_ALIGN.CENTER
                status_text = str(row.get("status", STATUS_ERROR_SYMBOL)).strip()
                is_ok = status_text in {STATUS_OK_SYMBOL, "V", "v", "✔"}
                color = STATUS_OK_GREEN if is_ok else STATUS_ERROR_RED
                paragraph.font.color.rgb = _to_rgb(rgb_color, color)
                if not paragraph.runs:
                    status_run = paragraph.add_run()
                    status_run.text = status_text
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.name = BODY_FONT_NAME
                    run.font.size = _scaled_pt(pt, body_font_size)
                    run.font.color.rgb = _to_rgb(rgb_color, color)


def _add_issue_slide(
    prs: Any,
    issue_index: int,
    issue_count: int,
    issue: dict[str, Any],
    rgb_color: Any,
    mso_anchor: Any,
    mso_auto_size: Any,
    inches: Any,
    pt: Any,
    issue_image_cache: dict[str, bytes | None],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    issue_id = str(issue.get("id", f"issue-{issue_index}"))
    wcag_rule = _display_wcag_rule(str(issue.get("wcag_rule", "N/A")))
    description = _normalize_text(issue.get("description"), fallback="N/A")
    source = _normalize_text(issue.get("source"), fallback="N/A")
    why_it_matters = _truncate(
        _normalize_text(issue.get("why_this_matters"), fallback="N/A"), WHY_IT_MATTERS_MAX_CHARS
    )
    potential_exposures = _format_potential_exposures(issue.get("potential_exposures"))
    snippet = _truncate(_normalize_text(issue.get("html_snippet"), fallback="N/A"), HTML_SNIPPET_MAX_CHARS)
    fix = _truncate(_normalize_text(issue.get("fix"), fallback="N/A"), FIX_MAX_CHARS)
    issue_image_source = _get_issue_image_source(issue)
    issue_image_bytes = _get_issue_image_bytes(issue_image_source, issue_image_cache)

    title_shape = slide.shapes.add_textbox(inches(0.4), inches(0.15), inches(9.2), inches(0.6))
    title_frame = title_shape.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = f"Issue {issue_index}/{issue_count}: {issue_id}"
    title_paragraph.font.size = _scaled_pt(pt, 21)
    title_paragraph.font.bold = True
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.color.rgb = _to_rgb(rgb_color, REPLY_GREEN)

    meta_shape = slide.shapes.add_textbox(inches(0.4), inches(0.85), inches(9.2), inches(2.05))
    meta_frame = meta_shape.text_frame
    meta_frame.clear()
    meta_frame.word_wrap = True
    meta_frame.auto_size = mso_auto_size.TEXT_TO_FIT_SHAPE

    meta_lines = [
        ("WCAG rule", wcag_rule),
        ("Description", description),
        ("Source", source),
        ("Why it matters", why_it_matters),
    ]
    for idx, (label, value) in enumerate(meta_lines):
        paragraph = meta_frame.paragraphs[0] if idx == 0 else meta_frame.add_paragraph()
        _set_label_value_paragraph(paragraph, label, value, pt, 12)

    snippet_left = inches(0.4)
    snippet_top = inches(2.5)
    snippet_height = inches(2.00)
    snippet_width = inches(9.2) if issue_image_bytes is None else inches(5.9)
    snippet_label_width = snippet_width

    snippet_label_shape = slide.shapes.add_textbox(inches(0.4), inches(2.93), snippet_label_width, inches(0.3))
    snippet_label_frame = snippet_label_shape.text_frame
    snippet_label_frame.clear()
    snippet_label_paragraph = snippet_label_frame.paragraphs[0]
    snippet_label_paragraph.text = "HTML snippet"
    snippet_label_paragraph.font.size = _scaled_pt(pt, 13)
    snippet_label_paragraph.font.bold = True
    snippet_label_paragraph.font.name = BODY_FONT_NAME

    snippet_shape = slide.shapes.add_textbox(snippet_left, snippet_top, snippet_width, snippet_height)
    snippet_shape.fill.solid()
    snippet_shape.fill.fore_color.rgb = rgb_color(236, 240, 241)
    snippet_shape.line.color.rgb = rgb_color(189, 195, 199)
    snippet_frame = snippet_shape.text_frame
    snippet_frame.clear()
    snippet_frame.word_wrap = True
    snippet_frame.vertical_anchor = mso_anchor.TOP
    snippet_paragraph = snippet_frame.paragraphs[0]
    snippet_paragraph.text = snippet
    snippet_paragraph.font.size = _scaled_pt(pt, 9)
    snippet_paragraph.font.name = SNIPPET_FONT_NAME

    if issue_image_bytes is not None:
        image_left = inches(6.45)
        image_top = inches(2.5)
        image_width = inches(3.15)
        image_height = inches(2.00)

        image_label_shape = slide.shapes.add_textbox(image_left, inches(2.93), image_width, inches(0.3))
        image_label_frame = image_label_shape.text_frame
        image_label_frame.clear()
        image_label_paragraph = image_label_frame.paragraphs[0]
        image_label_paragraph.text = "Issue image"
        image_label_paragraph.font.size = _scaled_pt(pt, 13)
        image_label_paragraph.font.bold = True
        image_label_paragraph.font.name = BODY_FONT_NAME

        image_box = slide.shapes.add_textbox(image_left, image_top, image_width, image_height)
        image_box.fill.solid()
        image_box.fill.fore_color.rgb = rgb_color(250, 250, 250)
        image_box.line.color.rgb = rgb_color(189, 195, 199)
        _add_centered_image_to_box(slide, issue_image_bytes, image_left, image_top, image_width, image_height)

    fix_label_shape = slide.shapes.add_textbox(inches(0.4), inches(4.9), inches(4.45), inches(0.3))
    fix_label_frame = fix_label_shape.text_frame
    fix_label_frame.clear()
    fix_label_paragraph = fix_label_frame.paragraphs[0]
    fix_label_paragraph.text = "Fix suggestion"
    fix_label_paragraph.font.size = _scaled_pt(pt, 13)
    fix_label_paragraph.font.bold = True
    fix_label_paragraph.font.name = BODY_FONT_NAME

    fix_shape = slide.shapes.add_textbox(inches(0.4), inches(5.2), inches(4.45), inches(1.85))
    fix_frame = fix_shape.text_frame
    fix_frame.clear()
    fix_frame.word_wrap = True
    fix_frame.auto_size = mso_auto_size.TEXT_TO_FIT_SHAPE
    fix_paragraph = fix_frame.paragraphs[0]
    fix_paragraph.text = fix
    fix_paragraph.font.size = _scaled_pt(pt, 11)
    fix_paragraph.font.name = BODY_FONT_NAME

    exposure_label_shape = slide.shapes.add_textbox(inches(5.05), inches(4.9), inches(4.55), inches(0.3))
    exposure_label_frame = exposure_label_shape.text_frame
    exposure_label_frame.clear()
    exposure_label_paragraph = exposure_label_frame.paragraphs[0]
    exposure_label_paragraph.text = "Potential exposures"
    exposure_label_paragraph.font.size = _scaled_pt(pt, 13)
    exposure_label_paragraph.font.bold = True
    exposure_label_paragraph.font.name = BODY_FONT_NAME

    exposure_shape = slide.shapes.add_textbox(inches(5.05), inches(5.2), inches(4.55), inches(1.85))
    exposure_frame = exposure_shape.text_frame
    exposure_frame.clear()
    exposure_frame.word_wrap = True
    exposure_frame.auto_size = mso_auto_size.TEXT_TO_FIT_SHAPE
    exposure_paragraph = exposure_frame.paragraphs[0]
    exposure_paragraph.text = potential_exposures
    exposure_paragraph.font.size = _scaled_pt(pt, 11)
    exposure_paragraph.font.name = BODY_FONT_NAME


def _get_a_aa_wcag_rules() -> list[str]:
    return list(dict.fromkeys(rule for rule in WCAG_RULE_MAPPER.values() if get_wcag_level(rule) in {"A", "AA"}))


def _level_order(level: object) -> int:
    if isinstance(level, str):
        return LEVEL_ORDER.get(level.upper(), 3)
    return 3


def _normalize_text(value: object, fallback: str) -> str:
    if value is None:
        return fallback
    text = str(value).strip()
    return text if text else fallback


def _format_potential_exposures(value: object) -> str:
    exposure_lines: list[str] = []
    for exposure in value if isinstance(value, list) else []:
        if not isinstance(exposure, dict):
            continue

        category = str(exposure.get("category") or "").strip()
        description = str(exposure.get("description") or "").strip()
        if category and description:
            exposure_lines.append(f"- {category}: {description}")
        elif description:
            exposure_lines.append(f"- {description}")

    if not exposure_lines:
        return "N/A"
    return _truncate("\n".join(exposure_lines), POTENTIAL_EXPOSURES_MAX_CHARS)


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return f"{text[: max_chars - 3]}..."


def _display_wcag_rule(value: str) -> str:
    if not value:
        return value
    return re.sub(r"\s*\(Level\s+[A-Z]+\)\s*$", "", value).strip()


def _to_rgb(rgb_color: Any, value: tuple[int, int, int]) -> Any:
    return rgb_color(value[0], value[1], value[2])


def _get_final_outcome(total_errors: int) -> tuple[str, tuple[int, int, int]]:
    if total_errors <= OUTCOME_FULLY_COMPLIANT_MAX_ERRORS:
        return OUTCOME_FULLY_COMPLIANT, STATUS_OK_GREEN
    if total_errors <= OUTCOME_PARTIALLY_COMPLIANT_MAX_ERRORS:
        return OUTCOME_PARTIALLY_COMPLIANT, OUTCOME_ORANGE
    return OUTCOME_NOT_COMPLIANT, STATUS_ERROR_RED


def _scaled_pt(pt: Any, font_size: int | float) -> Any:
    return pt(max(1, int(round(float(font_size) * FONT_SIZE_SCALE))))


def _set_label_value_paragraph(
    paragraph: Any,
    label: str,
    value: str,
    pt: Any,
    font_size: int,
    value_bold: bool = False,
    value_color: tuple[int, int, int] | None = None,
    rgb_color: Any | None = None,
) -> None:
    paragraph.clear()
    label_run = paragraph.add_run()
    label_run.text = f"{label}:"
    label_run.font.bold = True
    label_run.font.size = _scaled_pt(pt, font_size)
    label_run.font.name = BODY_FONT_NAME

    value_run = paragraph.add_run()
    value_run.text = f" {value}"
    value_run.font.bold = value_bold
    value_run.font.size = _scaled_pt(pt, font_size)
    value_run.font.name = BODY_FONT_NAME
    if value_color is not None and rgb_color is not None:
        value_run.font.color.rgb = _to_rgb(rgb_color, value_color)


def _chunk_rows(rows: list[dict[str, Any]], chunk_size: int) -> list[list[dict[str, Any]]]:
    if not rows:
        return [[]]
    return [rows[index : index + chunk_size] for index in range(0, len(rows), chunk_size)]


def _add_firstslide_template_slide(prs: Any, rgb_color: Any, inches: Any, pt: Any) -> None:
    _add_pdf_template_slide(
        prs=prs,
        filename=FIRSTSLIDE_TEMPLATE_FILENAME,
        fallback_title="Accessibility Report",
        rgb_color=rgb_color,
        inches=inches,
        pt=pt,
    )


def _add_thankyou_template_slide(prs: Any, rgb_color: Any, inches: Any, pt: Any) -> None:
    _add_pdf_template_slide(
        prs=prs,
        filename=THANKYOU_TEMPLATE_FILENAME,
        fallback_title="Thank you",
        rgb_color=rgb_color,
        inches=inches,
        pt=pt,
    )


def _add_pdf_template_slide(
    prs: Any,
    filename: str,
    fallback_title: str,
    rgb_color: Any,
    inches: Any,
    pt: Any,
) -> None:
    pdf_path = TEMPLATES_DIR / filename
    if not pdf_path.exists():
        return

    rendered_page = _render_first_pdf_page(pdf_path)
    if rendered_page is None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_template_fallback_text(slide, fallback_title, rgb_color, inches, pt)
        return

    image_bytes, image_width, image_height = rendered_page
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_centered_image_to_slide(
        slide=slide,
        image_bytes=image_bytes,
        image_width=image_width,
        image_height=image_height,
        slide_width=int(prs.slide_width),
        slide_height=int(prs.slide_height),
    )


def _add_logo_to_middle_slides(prs: Any, inches: Any) -> None:
    total_slides = len(prs.slides)
    if total_slides < 3:
        return

    for slide_index, slide in enumerate(prs.slides):
        if slide_index == 0 or slide_index == total_slides - 1:
            continue
        _add_logo_to_slide(prs, slide, inches)


def _add_logo_to_slide(prs: Any, slide: Any, inches: Any) -> None:
    logo_path = TEMPLATES_DIR / TEMPLATE_LOGO_FILENAME
    if not logo_path.exists():
        return

    logo_shape = slide.shapes.add_picture(str(logo_path), 0, 0, width=inches(LOGO_WIDTH_IN))
    margin = int(inches(LOGO_MARGIN_IN))
    logo_shape.left = int(prs.slide_width) - int(logo_shape.width) - margin
    logo_shape.top = int(prs.slide_height) - int(logo_shape.height) - margin


def _get_issue_image_source(issue: dict[str, Any]) -> str | None:
    raw_value = issue.get("img_url_or_path")
    if raw_value is None:
        raw_value = issue.get("image_url_or_path")
    if raw_value is None:
        return None

    normalized = str(raw_value).strip()
    if not normalized or normalized.lower() in {"none", "null"}:
        return None
    return normalized


def _get_issue_image_bytes(image_source: str | None, image_cache: dict[str, bytes | None]) -> bytes | None:
    if not image_source:
        return None
    if image_source in image_cache:
        return image_cache[image_source]

    try:
        if _is_http_url(image_source):
            image_bytes, mime = _download_issue_image_from_url(image_source)
        else:
            local_path = Path(image_source).expanduser()
            if not local_path.exists():
                image_cache[image_source] = None
                return None
            image_bytes = local_path.read_bytes()
            mime = _resolve_image_mime(image_source, "")

        if mime == "image/svg+xml" or image_source.lower().endswith(".svg"):
            image_bytes = _convert_svg_to_png(image_bytes)
        image_cache[image_source] = image_bytes
        return image_bytes
    except Exception:
        image_cache[image_source] = None
        return None


def _is_http_url(value: str) -> bool:
    parsed = urlparse(value)
    return parsed.scheme.lower() in {"http", "https"}


def _download_issue_image_from_url(url: str) -> tuple[bytes, str]:
    request = Request(url, headers={"User-Agent": "ax-tester/1.0"})
    with urlopen(request, timeout=ISSUE_IMAGE_DOWNLOAD_TIMEOUT) as response:
        content_type = response.headers.get("Content-Type", "")
        image_bytes = response.read(ISSUE_IMAGE_MAX_BYTES + 1)

    if len(image_bytes) > ISSUE_IMAGE_MAX_BYTES:
        raise ValueError(f"image_too_large>{ISSUE_IMAGE_MAX_BYTES}")
    return image_bytes, _resolve_image_mime(url, content_type)


def _resolve_image_mime(source: str, content_type: str) -> str:
    if content_type and "/" in content_type:
        return content_type.split(";")[0].strip()
    guessed, _ = mimetypes.guess_type(source)
    return guessed or "application/octet-stream"


def _convert_svg_to_png(svg_bytes: bytes) -> bytes:
    import cairosvg

    return cairosvg.svg2png(bytestring=svg_bytes)


def _add_centered_image_to_box(
    slide: Any,
    image_bytes: bytes,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
) -> None:
    try:
        picture = slide.shapes.add_picture(BytesIO(image_bytes), int(left), int(top), width=int(width))
    except Exception:
        return

    if int(picture.height) > int(height):
        picture.height = int(height)

    picture.left = int(left) + int((int(width) - int(picture.width)) / 2)
    picture.top = int(top) + int((int(height) - int(picture.height)) / 2)


def _render_first_pdf_page(pdf_path: Path) -> tuple[bytes, int, int] | None:
    try:
        import fitz
    except ImportError:
        return None

    with fitz.open(str(pdf_path)) as pdf_document:
        if pdf_document.page_count == 0:
            return None
        page = pdf_document.load_page(0)
        pixmap = page.get_pixmap(dpi=170, alpha=False)
        return pixmap.tobytes("png"), int(pixmap.width), int(pixmap.height)


def _add_centered_image_to_slide(
    slide: Any,
    image_bytes: bytes,
    image_width: int,
    image_height: int,
    slide_width: int,
    slide_height: int,
) -> None:
    if image_width <= 0 or image_height <= 0:
        return

    scale = min(slide_width / image_width, slide_height / image_height)
    rendered_width = max(1, int(image_width * scale))
    rendered_height = max(1, int(image_height * scale))
    left = int((slide_width - rendered_width) / 2)
    top = int((slide_height - rendered_height) / 2)
    slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=rendered_width, height=rendered_height)


def _add_template_fallback_text(slide: Any, title_text: str, rgb_color: Any, inches: Any, pt: Any) -> None:
    title_shape = slide.shapes.add_textbox(inches(0.5), inches(2.2), inches(9.0), inches(1.0))
    title_frame = title_shape.text_frame
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title_text
    title_paragraph.font.size = _scaled_pt(pt, 54)
    title_paragraph.font.bold = True
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.color.rgb = _to_rgb(rgb_color, REPLY_GREEN)


if __name__ == "__main__":
    results_dir = "/home/pbianco/ax_tester/results/2026-05-18_15-54-29_shop.reply.com/2026-05-18_15-58-56"
    build_pptx_report(results_dir)
